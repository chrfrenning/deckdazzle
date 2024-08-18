from pptx import Presentation
from pptx.util import Inches
from unsplash import download_random_image_by_keyword

# Global variable to hold the presentation object
prs = Presentation("template.pptx")

def createTitleSlide(title_text):
    """Create a title slide with the given title and subtitle."""
    slide_layout = prs.slide_layouts[0]  # 0 is the layout for title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text

def createBulletSlide(title_text, bullet_points):
    """Create a bullet slide with the given title and bullet points."""
    slide_layout = prs.slide_layouts[1]  # 1 is the layout for title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    content = slide.placeholders[1].text_frame
    for point in bullet_points:
        p = content.add_paragraph()
        p.text = point

def createTitleAndImageSlide(title_text, image_path):
    """Create a slide with the given title and image."""
    slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    slide.placeholders[1].insert_picture(image_path)

def createPictureSlide(image_path):
    """Create a picture slide with the given image only."""
    slide_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(slide_layout)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)

def save(file_path):
    """Save the presentation to the given file path."""
    prs.save(file_path)


def main():
    # Call the functions to create slides
    createTitleSlide("Welcome to Deck Dazzle")
    
    bullet_points = ["Point 1", "Point 2", "Point 3"]
    createBulletSlide("Bullet Slide", bullet_points)

    download_random_image_by_keyword(["twin", "transformation"], "unsplash.jpg")
    createTitleAndImageSlide("Title and Image Slide", "unsplash.jpg")
    
    createPictureSlide("unsplash.jpg")
    
    # Save the presentation
    save("sample_presentation.pptx")

if __name__ == "__main__":
    main()