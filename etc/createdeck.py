from pptx import Presentation
from pptx.util import Inches

# Create a Presentation object
template_path = 'template.pptx'  # Replace with the path to your template
prs = Presentation(template_path)

# Add a title slide
slide_layout = prs.slide_layouts[0]  # 0 is the layout for title slide
slide = prs.slides.add_slide(slide_layout)

# Define title and subtitle for the slide
title = slide.shapes.title

title.text = "My Presentation Title"

# Add a slide with title and content
slide_layout = prs.slide_layouts[1]  # 1 is the layout for title and content
slide = prs.slides.add_slide(slide_layout)

# Define title for the content slide
title = slide.shapes.title
title.text = "Second Slide Title"

# Add content to the slide
content = slide.shapes.placeholders[1]
content.text = "Here is some content for the second slide."

# Add a bullet point list
bullet_points = [
    "First bullet point",
    "Second bullet point",
    "Third bullet point"
]

for point in bullet_points:
    p = content.text_frame.add_paragraph()
    p.text = point
    p.level = 0  # level 0 is the top level


# create a slide with a full page image

slide_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(slide_layout)

# Define the path to your image
img_path = 'unsplash.jpg'  # Replace with the path to your image

# Get slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

# Add image to slide, making it fill the entire slide
slide.shapes.add_picture(img_path, 0, 0, width=slide_width, height=slide_height)


# Save the presentation
prs.save('sample_presentation.pptx')

print("Presentation created successfully.")